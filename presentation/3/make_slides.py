"""Generate summary slides for the NNCL2 project."""

import io
import pathlib
import subprocess
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (light theme) ─────────────────────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white
ACCENT  = RGBColor(0x1A, 0x5C, 0xC8)   # strong blue
TEXT    = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
SUBTEXT = RGBColor(0x55, 0x55, 0x6E)   # muted
BOX_BG  = RGBColor(0xF0, 0xF3, 0xFA)   # very light blue-gray
RULE    = RGBColor(0xD0, 0xD8, 0xF0)   # faint divider
W = Inches(13.33)
H = Inches(7.5)

HEADER_H = Inches(1.0)

# Hex strings for matplotlib (it doesn't take RGBColor)
_ACCENT_HEX  = "#1A5CC8"
_BOX_BG_HEX  = "#F0F3FA"


# ── Equation renderer ─────────────────────────────────────────────────────────

_EQ_DPI      = 300   # render resolution
_EQ_FONTSIZE = 13    # pt in the rendered image


def render_eq(latex: str) -> io.BytesIO:
    """Render a LaTeX math string to a high-res PNG BytesIO using matplotlib mathtext."""
    fig = plt.figure(figsize=(0.01, 0.01), dpi=_EQ_DPI)
    fig.patch.set_facecolor(_BOX_BG_HEX)
    text = fig.text(0, 0, f"${latex}$",
                    fontsize=_EQ_FONTSIZE, color=_ACCENT_HEX,
                    ha="left", va="bottom")
    fig.canvas.draw()
    bbox = text.get_window_extent(renderer=fig.canvas.get_renderer())
    fig.set_size_inches((bbox.width + 4) / _EQ_DPI, (bbox.height + 4) / _EQ_DPI)
    text.set_position((2 / bbox.width, 2 / bbox.height))
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=_EQ_DPI,
                facecolor=_BOX_BG_HEX, edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def add_eq(slide, latex: str, left, top, width=Inches(2.4), height=None):
    """Render equation and insert it. height takes priority over width if given."""
    buf = render_eq(latex)
    if height is not None:
        slide.shapes.add_picture(buf, left, top, height=height)
    else:
        slide.shapes.add_picture(buf, left, top, width=width)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def draw_bullets(slide, bullets, left, top, width, blh,
                 size=11, color=TEXT, indent=Inches(0.18), line_spacing=1.2):
    """Single textbox, all bullets in one paragraph separated by line breaks."""
    text = "\n".join("•  " + b for b in bullets)
    n = len(bullets)
    tb = add_text(slide, text, left + indent, top, width - 2 * indent, blh * n,
                  size=size, color=color, wrap=True)
    tb.text_frame.paragraphs[0].line_spacing = line_spacing
    return top + blh * n


def page_header(slide, title):
    """Accent bar + title at the top."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=ACCENT)
    add_text(slide, title,
             Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
             size=28, bold=True, color=ACCENT)
    add_rect(slide, Inches(0.5), HEADER_H - Inches(0.04), W - Inches(1.0), Inches(0.02),
             fill_color=RULE)


def card(slide, left, top, width, height, title=None, items=None,
         item_size=15, title_size=16, item_color=TEXT):
    """Light card with optional title and bullet items."""
    add_rect(slide, left, top, width, height, fill_color=BOX_BG)
    y = top + Inches(0.18)
    if title:
        add_text(slide, title, left + Inches(0.2), y, width - Inches(0.4), Inches(0.38),
                 size=title_size, bold=True, color=ACCENT)
        y += Inches(0.38)
    if items:
        text = "\n".join(f"  {item}" for item in items)
        add_text(slide, text, left + Inches(0.15), y,
                 width - Inches(0.3), height - (y - top) - Inches(0.1),
                 size=item_size, color=item_color, wrap=True)


# ── Slide 1 – Objectives ──────────────────────────────────────────────────────

def slide_objectives(prs):
    slide = blank_slide(prs)
    set_bg(slide, BG)
    page_header(slide, "Objectives")

    add_text(slide,
             "We study the internal representations of a SimCLR model — "
             "focusing on the geometric and structural properties of the learned features.",
             Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.6),
             size=16, color=SUBTEXT)

    # (label, what, why)
    questions = [
        ("Sparsity",
         "Do most feature dimensions carry little signal?",
         "Sparse representations are more interpretable and may indicate cleaner feature reuse."),
        ("Separability",
         "How well does each dimension separate classes?",
         "High separability means features are directly useful for classification — a proxy for representation quality."),
        ("Disentanglement",
         "Is each dimension selective for a single class?",
         "Disentangled features are easier to interpret and more robust to distribution shift."),
        ("Completeness",
         "Is each class explained by a small set of dimensions?",
         "Low completeness means a class is spread across many dims, making its representation harder to interpret."),
        ("Layer comparison",
         "How do these properties differ across encoder output and projector?",
         "The projector is discarded at inference — understanding how it transforms features explains why SSL works."),
    ]

    col_w = Inches(5.9)
    gap   = Inches(0.43)
    y0      = Inches(1.75)
    row_gap = Inches(0.14)
    # 2 full rows + 1 full-width row at the bottom
    avail = H - Inches(0.25) - y0
    row_h = (avail - 2 * row_gap) / 3

    def draw_card(lft, top, width, label, what, why):
        add_rect(slide, lft, top, width, row_h, fill_color=BOX_BG)
        tb = slide.shapes.add_textbox(
            lft + Inches(0.18), top + Inches(0.1),
            width - Inches(0.36), row_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (text, size, bold, clr) in enumerate([
            (label, 14, True,  ACCENT),
            (what,  13, False, TEXT),
            (why,   12, False, SUBTEXT),
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = clr

    # Rows 0 and 1: two columns each
    for i in range(4):
        col = i % 2
        row = i // 2
        lft = Inches(0.5) + col * (col_w + gap)
        top = y0 + row * (row_h + row_gap)
        draw_card(lft, top, col_w, *questions[i])

    # Row 2: last item spans full width
    top = y0 + 2 * (row_h + row_gap)
    draw_card(Inches(0.5), top, col_w * 2 + gap, *questions[4])


# ── Slide 2 – Setup ───────────────────────────────────────────────────────────

def slide_setup(prs):
    slide = blank_slide(prs)
    set_bg(slide, BG)
    page_header(slide, "Experimental Setup")

    top = Inches(1.1)
    col_w = Inches(3.9)
    gap   = Inches(0.27)
    h     = Inches(5.6)

    # Dataset
    card(slide,
         left=Inches(0.5), top=top, width=col_w, height=h,
         title="Dataset — CIFAR-100",
         items=[
             "50 000 train / 10 000 test",
             "100 classes, 32 × 32 px",
             "",
             "SimCLR augmentation:",
             "  · Random crop + h-flip",
             "  · Colour jitter",
             "  · Gaussian blur",
             "  · Random grayscale",
             "",
             "Two augmented views",
             " · Positive pair in constrastive learning",
         ],
         item_size=14)

    # Architecture
    card(slide,
         left=Inches(0.5) + col_w + gap, top=top, width=col_w, height=h,
         title="Architecture",
         items=[
             "Encoder: ResNet-18",
             "  · 3×3 stride-1 first conv",
             "    (preserves 32px resolution)",
             "  · Output: 512-d vector",
             "",
             "Projector (optional)",
             "  · Linear",
             "  · BatchNorm (optional)",
             "  · ReLU",
             "  · Linear",
             "  · BatchNorm (optional)",
         ],
         item_size=14)

    # Training
    card(slide,
         left=Inches(0.5) + 2*(col_w + gap), top=top, width=col_w, height=h,
         title="Training",
         items=[
             "Loss: NT-Xent  (Temperature = 0.1)",
             "  · 2(B-1) negatives",
             "",
             "Optimiser: SGD",
             "  · Momentum 0.9",
             "  · Weight decay 1e-4",
             "  · (excl. BN & biases)",
             "",
             "Schedule:",
             "  · Warmup 10 ep -> cosine annealing",
             "  · 256 batch size"
             "  · 200 epochs",
         ],
         item_size=14)


# ── Slide 3 – Metrics ─────────────────────────────────────────────────────────

def slide_metrics(prs):
    slide = blank_slide(prs)
    set_bg(slide, BG)
    page_header(slide, "Metrics")

    PAD        = Inches(0.5)
    INNER      = Inches(0.18)
    CW         = W - 2 * PAD
    BLH        = Inches(0.27)   # fixed bullet line height — pack from top, no stretching

    def section_label(text, y):
        add_text(slide, text, PAD, y, CW, Inches(0.3),
                 size=12, bold=True, color=SUBTEXT)
        add_rect(slide, PAD, y + Inches(0.28), CW, Inches(0.015), fill_color=RULE)

    def metric_card(left, top, width, height, title, bullets, eq_latex=None):
        add_rect(slide, left, top, width, height, fill_color=BOX_BG)
        add_text(slide, title,
                 left + INNER, top + Inches(0.1), width - 2*INNER, Inches(0.28),
                 size=13, bold=True, color=ACCENT)
        y = draw_bullets(slide, bullets, left, top + Inches(0.42), width, BLH)
        if eq_latex:
            add_eq(slide, eq_latex, left + INNER, y + Inches(0.08))

    # ── Downstream ────────────────────────────────────────────────────────────
    section_label("Downstream Accuracy", Inches(1.05))

    d_top = Inches(1.42)
    d_h   = Inches(1.55)
    d_cw  = (CW - Inches(0.3)) / 2

    metric_card(PAD, d_top, d_cw, d_h,
        "kNN Accuracy  (k = 20)",
        ["0 = random; 1 = perfect - tests geometric clustering of features",
         "Majority class among 20 nearest neighbours",
         "No training, no linear assumption - pure geometry",
         "Reports Top-1 and Top-5"])

    metric_card(PAD + d_cw + Inches(0.3), d_top, d_cw, d_h,
        "Linear Probe  (Top-1 / Top-5)",
        ["Tests whether class info is linearly readable from frozen features",
         "Single linear layer on encoder output",
         "Good accuracy = no non-linear decoding needed",
         "Reports Top-1 and Top-5"])

    # ── Sparsity ──────────────────────────────────────────────────────────────
    s_top0 = d_top + d_h + Inches(0.16)
    section_label("Sparsity Metrics", s_top0)

    s_top = s_top0 + Inches(0.36)
    avail = H - Inches(0.2) - s_top
    s_h   = avail
    s_gap = Inches(0.22)
    s_cw  = (CW - 2 * s_gap) / 3

    def sx(col): return PAD + col * (s_cw + s_gap)

    metric_card(sx(0), s_top, s_cw, s_h,
        "Activation histograms",
        ["Tracks the distribution of activations",
         "Overview over the feature geometry"])

    metric_card(sx(1), s_top, s_cw, s_h,
        "Zero Fraction",
        ["Fraction of activations that are very small",
         "High = many neurons off for a given input",
         "Measures input-specific hard sparsity"])

    metric_card(sx(2), s_top, s_cw, s_h,
        "Hoyer Sparsity Index",
        ["0 = uniform activity; 1 = single active unit",
         "L1/L2 ratio - normalised to [0, 1]"],
        eq_latex=r"\mathrm{HSI} = \frac{\sqrt{D} - \|x\|_1\,/\,\|x\|_2}{\sqrt{D} - 1} \in [0,\,1]")


# ── Slide 4 – Representation Structure ───────────────────────────────────────

def slide_structure(prs):
    slide = blank_slide(prs)
    set_bg(slide, BG)
    page_header(slide, "Representation Structure Metrics")

    PAD   = Inches(0.5)
    INNER = Inches(0.18)
    CW    = W - 2 * PAD
    BLH   = Inches(0.27)
    EQ_H  = Inches(0.42)

    # ── Intro box: intuition behind the variance ratio ────────────────────────
    ib_top = Inches(1.38)
    ib_h   = Inches(1.75)
    ib_lw  = CW * 0.52   # left text column
    ib_rx  = PAD + ib_lw + Inches(0.3)

    add_rect(slide, PAD, ib_top, CW, ib_h, fill_color=BOX_BG)

    # Left: verbal intuition
    intro_bullets = [
        "Low within-class variance indicates tighly clustered representations",
        "High between-class variance indicates separated representations",
        "The ratio between / (within + between) quantifies the amount of variability explained by class difference (bounded [0,1])",
        "=> Variance decomposition"
    ]
    draw_bullets(slide, intro_bullets, PAD, ib_top + Inches(0.14), ib_lw, BLH)

    # Right: the key equations stacked
    eq_y = ib_top + Inches(0.12)
    add_eq(slide,
           r"\sigma^2_W = \sum_c w_c\,\sigma^2_{c,d}",
           ib_rx, eq_y, height=Inches(0.32))
    eq_y += Inches(0.38)
    add_eq(slide,
           r"\sigma^2_B = \sum_c w_c\,(\mu_{c,d}-\mu_d)^2",
           ib_rx, eq_y, height=Inches(0.32))
    eq_y += Inches(0.38)
    add_eq(slide,
           r"s_d = \frac{\sigma^2_B}{\sigma^2_W + \sigma^2_B} \in [0,\,1]",
           ib_rx, eq_y, height=Inches(0.42))

    # ── Three metric cards ────────────────────────────────────────────────────
    s_gap = Inches(0.22)
    s_cw  = (CW - 2 * s_gap) / 3
    s_top = ib_top + ib_h + Inches(0.16)
    s_h   = H - Inches(0.2) - s_top

    def sx(col): return PAD + col * (s_cw + s_gap)

    def struct_card(left, top, width, height, title, eq_latex, r_eq, bullets):
        add_rect(slide, left, top, width, height, fill_color=BOX_BG)
        add_text(slide, title,
                 left + INNER, top + Inches(0.1), width - 2*INNER, Inches(0.28),
                 size=13, bold=True, color=ACCENT)
        y = top + Inches(0.42)
        if r_eq:
            add_eq(slide, r_eq, left + INNER, y, height=EQ_H)
            y += EQ_H + Inches(0.06)
        add_eq(slide, eq_latex, left + INNER, y, height=EQ_H)
        y += EQ_H + Inches(0.06)
        draw_bullets(slide, bullets, left, y, width, BLH)

    # Separability
    struct_card(sx(0), s_top, s_cw, s_h,
        "Separability",
        eq_latex=r"s_d = \frac{\sigma^2_B}{\sigma^2_W + \sigma^2_B} \in [0,\,1]",
        r_eq=None,
        bullets=[
            "How strongly does a feature predict class identity?",
            "0 = no class signal; 1 = pure class signal",
            "Averaged over all features for a scalar score",
        ])

    # Disentanglement
    struct_card(sx(1), s_top, s_cw, s_h,
        "Disentanglement",
        eq_latex=r"D_i = 1 - \frac{H(P_{i.})}{\log K}",
        r_eq=r"P_{ij} = \frac{R_{ij}}{\sum_k R_{ik}}",
        bullets=[
            "Does each latent feature only represent one class?",
            "R: (Latent x Class) matrix of relative importance",
            "R_ij: importance of feature i for class j",
            "Here: Variance explained ratio as importance for every feature/class",
            "Simple terms: 1 - per-feature entropy of normalized relative importances",
            "Low entropy => Feature variance only explained by a single class => High Disentanglement",
        ])

    # Completeness
    struct_card(sx(2), s_top, s_cw, s_h,
        "Completeness",
        eq_latex=r"C_j = 1 - \frac{H(\tilde{P}_{.j})}{\log D}",
        r_eq=r"\tilde{P}_{ij} = \frac{R_{ij}}{\sum_k R_{kj}}",
        bullets=[
            "Is each class only explained by a single latent feature?",
            "Mirrored to disentanglement: 1 - entropy of per-class (normalized) relative importances",
            "Low entropy => Only a single feature indicates this class identity well => High completeness",
        ])

    # ── Paper reference ───────────────────────────────────────────────────────
    add_text(slide,
             "Eastwood & Williams, ICLR 2018  —  "
             "\"A Framework for the Quantitative Evaluation of Disentangled Representations\"",
             PAD, Inches(1.05), CW, Inches(0.26), size=11, color=SUBTEXT)


# ── Build & save ──────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_setup(prs)
    slide_metrics(prs)
    slide_structure(prs)

    out = pathlib.Path(__file__).parent / "slides.pptx"
    prs.save(out)
    print(f"Saved → {out}")

    lookup = "where" if sys.platform == "win32" else "which"
    lo = next(
        (p for p in ["soffice", "libreoffice",
                     "/opt/homebrew/bin/soffice",
                     "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
         if subprocess.run([lookup, p], capture_output=True).returncode == 0
         or pathlib.Path(p).exists()),
        None,
    )
    if lo is None:
        print("LibreOffice not found — skipping PDF conversion")
        return
    result = subprocess.run(
        [lo, "--headless", "--invisible", "--norestore", "--nofirststartwizard",
         "--convert-to", "pdf", "--outdir", str(out.parent), str(out)],
        capture_output=True, text=True,
    )
    if result.returncode == 0:
        print(f"PDF  → {out.with_suffix('.pdf')}")
    else:
        print(f"LibreOffice failed:\n{result.stderr}")


if __name__ == "__main__":
    main()
